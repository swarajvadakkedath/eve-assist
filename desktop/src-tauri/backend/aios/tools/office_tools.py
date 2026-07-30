"""Office Document Toolkit — PDF, DOCX, XLSX, PPTX for AIOS Phase 5.4B."""

from __future__ import annotations

import asyncio
import os
from pathlib import Path
from typing import Any

from aios.core.tool_manager import ToolResult
from aios.core.event_bus import EventBus


# ── PDF Tools ──


async def _read_pdf(params: dict) -> ToolResult:
    try:
        import fitz

        path = Path(params["path"])
        if not path.exists():
            return ToolResult(success=False, error=f"File not found: {path}")

        doc = fitz.open(str(path))
        pages = []
        for i in range(len(doc)):
            page = doc[i]
            pages.append({
                "page": i + 1,
                "text": page.get_text().strip(),
                "char_count": len(page.get_text()),
            })

        metadata = doc.metadata or {}

        return ToolResult(success=True, data={
            "path": str(path),
            "pages": pages,
            "page_count": len(pages),
            "metadata": {
                "title": metadata.get("title", ""),
                "author": metadata.get("author", ""),
                "subject": metadata.get("subject", ""),
                "format": metadata.get("format", ""),
                "producer": metadata.get("producer", ""),
                "creator": metadata.get("creator", ""),
            },
        })
    except ImportError:
        return ToolResult(success=False, error="PyMuPDF is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


async def _search_pdf(params: dict) -> ToolResult:
    try:
        import fitz

        path = Path(params["path"])
        query = params.get("query", "")
        case_sensitive = params.get("case_sensitive", False)

        if not path.exists():
            return ToolResult(success=False, error=f"File not found: {path}")
        if not query:
            return ToolResult(success=False, error="No search query provided")

        doc = fitz.open(str(path))
        results = []
        q = query if case_sensitive else query.lower()

        for i in range(len(doc)):
            page = doc[i]
            text = page.get_text()
            check = text if case_sensitive else text.lower()
            if q not in check:
                continue

            lines = text.splitlines()
            for j, line in enumerate(lines):
                cline = line if case_sensitive else line.lower()
                if q in cline:
                    results.append({
                        "page": i + 1,
                        "line": j + 1,
                        "content": line.strip()[:500],
                    })

        return ToolResult(success=True, data={
            "path": str(path),
            "query": query,
            "results": results,
            "matches": len(results),
        })
    except ImportError:
        return ToolResult(success=False, error="PyMuPDF is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


async def _extract_pdf_metadata(params: dict) -> ToolResult:
    try:
        import fitz

        path = Path(params["path"])
        if not path.exists():
            return ToolResult(success=False, error=f"File not found: {path}")

        doc = fitz.open(str(path))
        metadata = doc.metadata or {}
        toc = doc.get_toc()

        return ToolResult(success=True, data={
            "path": str(path),
            "page_count": len(doc),
            "file_size": path.stat().st_size,
            "metadata": {
                "title": metadata.get("title", ""),
                "author": metadata.get("author", ""),
                "subject": metadata.get("subject", ""),
                "keywords": metadata.get("keywords", ""),
                "producer": metadata.get("producer", ""),
                "creator": metadata.get("creator", ""),
                "format": metadata.get("format", ""),
                "creation_date": metadata.get("creationDate", ""),
                "mod_date": metadata.get("modDate", ""),
            },
            "has_toc": len(toc) > 0,
            "toc_entries": len(toc),
        })
    except ImportError:
        return ToolResult(success=False, error="PyMuPDF is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


async def _split_pdf(params: dict, event_bus: EventBus | None = None) -> ToolResult:
    try:
        import fitz

        path = Path(params["path"])
        output_dir = Path(params.get("output_dir", ""))
        pages = params.get("pages", [])

        if not path.exists():
            return ToolResult(success=False, error=f"File not found: {path}")

        doc = fitz.open(str(path))

        if not output_dir:
            output_dir = path.parent / f"{path.stem}_split"
        output_dir.mkdir(parents=True, exist_ok=True)

        result_files = []

        if pages:
            for page_num in pages:
                if page_num < 1 or page_num > len(doc):
                    continue
                out_path = output_dir / f"{path.stem}_page_{page_num}.pdf"
                new_doc = fitz.open()
                new_doc.insert_pdf(doc, from_page=page_num - 1, to_page=page_num - 1)
                new_doc.save(str(out_path))
                new_doc.close()
                result_files.append({"page": page_num, "path": str(out_path)})
        else:
            for i in range(len(doc)):
                out_path = output_dir / f"{path.stem}_page_{i + 1}.pdf"
                new_doc = fitz.open()
                new_doc.insert_pdf(doc, from_page=i, to_page=i)
                new_doc.save(str(out_path))
                new_doc.close()
                result_files.append({"page": i + 1, "path": str(out_path)})

        doc.close()

        if event_bus:
            await event_bus.publish(
                "office:pdf_split",
                {"source": str(path), "output_dir": str(output_dir), "files": len(result_files)},
                source="office_tools",
            )

        return ToolResult(success=True, data={
            "source": str(path),
            "output_dir": str(output_dir),
            "files": result_files,
            "count": len(result_files),
        })
    except ImportError:
        return ToolResult(success=False, error="PyMuPDF is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


async def _merge_pdf(params: dict, event_bus: EventBus | None = None) -> ToolResult:
    try:
        import fitz

        sources = params.get("sources", [])
        output = Path(params.get("output", ""))

        if not sources:
            return ToolResult(success=False, error="No source files provided")
        if not output:
            return ToolResult(success=False, error="No output path provided")

        output.parent.mkdir(parents=True, exist_ok=True)
        merged = fitz.open()
        file_results = []

        for src in sources:
            src_path = Path(src)
            if not src_path.exists():
                file_results.append({"file": src, "error": "File not found"})
                continue
            try:
                doc = fitz.open(str(src_path))
                merged.insert_pdf(doc)
                file_results.append({"file": src, "pages": len(doc)})
                doc.close()
            except Exception as e:
                file_results.append({"file": src, "error": str(e)})

        total_pages = merged.page_count
        merged.save(str(output))
        merged.close()

        if event_bus:
            await event_bus.publish(
                "office:pdf_merge",
                {"sources": sources, "output": str(output), "total_pages": total_pages},
                source="office_tools",
            )

        return ToolResult(success=True, data={
            "output": str(output),
            "files": file_results,
            "total_files": len(sources),
        })
    except ImportError:
        return ToolResult(success=False, error="PyMuPDF is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


# ── DOCX Tools ──


async def _read_docx(params: dict) -> ToolResult:
    try:
        from docx import Document

        path = Path(params["path"])
        if not path.exists():
            return ToolResult(success=False, error=f"File not found: {path}")

        doc = Document(str(path))
        paragraphs = []
        for p in doc.paragraphs:
            style = p.style.name if p.style else "Normal"
            paragraphs.append({
                "text": p.text,
                "style": style,
                "is_heading": "Heading" in style,
                "heading_level": int(style.split()[-1]) if "Heading" in style and style.split()[-1].isdigit() else 0,
            })

        return ToolResult(success=True, data={
            "path": str(path),
            "paragraphs": paragraphs,
            "paragraph_count": len(paragraphs),
            "headings": [p for p in paragraphs if p["is_heading"]],
        })
    except ImportError:
        return ToolResult(success=False, error="python-docx is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


async def _write_docx(params: dict, event_bus: EventBus | None = None) -> ToolResult:
    try:
        from docx import Document

        path = Path(params["path"])
        content = params.get("content", [])

        path.parent.mkdir(parents=True, exist_ok=True)
        doc = Document()

        for item in content:
            text = item.get("text", "")
            style = item.get("style", "paragraph")

            if style == "title":
                doc.add_heading(text, level=0)
            elif style == "heading1":
                doc.add_heading(text, level=1)
            elif style == "heading2":
                doc.add_heading(text, level=2)
            elif style == "heading3":
                doc.add_heading(text, level=3)
            else:
                doc.add_paragraph(text)

        doc.save(str(path))

        if event_bus:
            await event_bus.publish(
                "office:docx_write",
                {"path": str(path), "items": len(content)},
                source="office_tools",
            )

        return ToolResult(success=True, data={
            "path": str(path),
            "items_written": len(content),
        })
    except ImportError:
        return ToolResult(success=False, error="python-docx is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


async def _extract_headings(params: dict) -> ToolResult:
    try:
        from docx import Document

        path = Path(params["path"])
        if not path.exists():
            return ToolResult(success=False, error=f"File not found: {path}")

        doc = Document(str(path))
        headings = []
        for p in doc.paragraphs:
            if p.style and "Heading" in p.style.name:
                level = int(p.style.name.split()[-1]) if p.style.name.split()[-1].isdigit() else 1
                headings.append({
                    "text": p.text,
                    "level": level,
                    "style": p.style.name,
                })

        return ToolResult(success=True, data={
            "path": str(path),
            "headings": headings,
            "count": len(headings),
        })
    except ImportError:
        return ToolResult(success=False, error="python-docx is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


# ── XLSX Tools ──


async def _list_sheets(params: dict) -> ToolResult:
    try:
        from openpyxl import load_workbook

        path = Path(params["path"])
        if not path.exists():
            return ToolResult(success=False, error=f"File not found: {path}")

        wb = load_workbook(str(path), read_only=True)
        sheets = []
        for name in wb.sheetnames:
            ws = wb[name]
            sheets.append({
                "name": name,
                "rows": ws.max_row or 0,
                "columns": ws.max_column or 0,
            })
        wb.close()

        return ToolResult(success=True, data={
            "path": str(path),
            "sheets": sheets,
            "count": len(sheets),
        })
    except ImportError:
        return ToolResult(success=False, error="openpyxl is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


async def _read_sheet(params: dict) -> ToolResult:
    try:
        from openpyxl import load_workbook

        path = Path(params["path"])
        sheet_name = params.get("sheet", "")
        max_rows = params.get("max_rows", 10000)
        has_header = params.get("has_header", True)

        if not path.exists():
            return ToolResult(success=False, error=f"File not found: {path}")

        wb = load_workbook(str(path), read_only=True)

        if sheet_name:
            if sheet_name not in wb.sheetnames:
                wb.close()
                return ToolResult(success=False, error=f"Sheet not found: {sheet_name}")
            ws = wb[sheet_name]
        else:
            ws = wb.active

        rows = []
        headers = None
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= max_rows:
                break
            cleaned = [_clean_cell(v) for v in row]
            if i == 0 and has_header:
                headers = [str(v) if v is not None else f"Col_{j}" for j, v in enumerate(cleaned)]
            else:
                if headers:
                    rows.append(dict(zip(headers, cleaned)))
                else:
                    rows.append(cleaned)

        wb.close()

        return ToolResult(success=True, data={
            "path": str(path),
            "sheet": sheet_name or (wb.sheetnames[0] if wb.sheetnames else "unknown"),
            "headers": headers,
            "rows": rows,
            "row_count": len(rows),
        })
    except ImportError:
        return ToolResult(success=False, error="openpyxl is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


def _clean_cell(value: Any) -> Any:
    if value is None:
        return None
    if isinstance(value, (int, float)):
        return value
    return str(value)


async def _write_sheet(params: dict, event_bus: EventBus | None = None) -> ToolResult:
    try:
        from openpyxl import Workbook

        path = Path(params["path"])
        sheet_name = params.get("sheet", "Sheet1")
        data = params.get("data", [])
        headers = params.get("headers", [])

        if not data:
            return ToolResult(success=False, error="No data provided")

        path.parent.mkdir(parents=True, exist_ok=True)
        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name

        if headers:
            ws.append(headers)

        for row in data:
            if isinstance(row, dict):
                if headers:
                    ws.append([row.get(h, "") for h in headers])
                else:
                    ws.append(list(row.values()))
            elif isinstance(row, (list, tuple)):
                ws.append(list(row))
            else:
                ws.append([row])

        wb.save(str(path))

        if event_bus:
            await event_bus.publish(
                "office:xlsx_write",
                {"path": str(path), "sheet": sheet_name, "rows": len(data)},
                source="office_tools",
            )

        return ToolResult(success=True, data={
            "path": str(path),
            "sheet": sheet_name,
            "rows_written": len(data),
            "headers": headers,
        })
    except ImportError:
        return ToolResult(success=False, error="openpyxl is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


# ── PPTX Tools ──


async def _read_presentation(params: dict) -> ToolResult:
    try:
        from pptx import Presentation

        path = Path(params["path"])
        if not path.exists():
            return ToolResult(success=False, error=f"File not found: {path}")

        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides):
            shapes = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        shapes.append({
                            "type": str(shape.shape_type),
                            "text": text[:1000],
                            "name": shape.name,
                        })
            slide_data = {
                "slide": i + 1,
                "shapes": shapes,
                "shape_count": len(shapes),
            }
            title_shape = slide.shapes.title
            slide_data["title"] = title_shape.text if title_shape else ""
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                slide_data["notes"] = notes if notes else ""
            slides.append(slide_data)

        return ToolResult(success=True, data={
            "path": str(path),
            "slides": slides,
            "slide_count": len(slides),
        })
    except ImportError:
        return ToolResult(success=False, error="python-pptx is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


async def _list_slides(params: dict) -> ToolResult:
    try:
        from pptx import Presentation

        path = Path(params["path"])
        if not path.exists():
            return ToolResult(success=False, error=f"File not found: {path}")

        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides):
            title_shape = slide.shapes.title
            title = title_shape.text if title_shape else ""
            slides.append({
                "slide": i + 1,
                "title": title,
                "has_notes": slide.has_notes_slide,
                "shape_count": len(slide.shapes),
            })

        return ToolResult(success=True, data={
            "path": str(path),
            "slides": slides,
            "count": len(slides),
        })
    except ImportError:
        return ToolResult(success=False, error="python-pptx is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


async def _extract_notes(params: dict) -> ToolResult:
    try:
        from pptx import Presentation

        path = Path(params["path"])
        if not path.exists():
            return ToolResult(success=False, error=f"File not found: {path}")

        prs = Presentation(str(path))
        notes_list = []
        for i, slide in enumerate(prs.slides):
            note_text = ""
            if slide.has_notes_slide:
                note_text = slide.notes_slide.notes_text_frame.text.strip()
            notes_list.append({
                "slide": i + 1,
                "has_notes": bool(note_text),
                "notes": note_text,
            })

        return ToolResult(success=True, data={
            "path": str(path),
            "notes": notes_list,
            "total_slides": len(notes_list),
            "slides_with_notes": sum(1 for n in notes_list if n["has_notes"]),
        })
    except ImportError:
        return ToolResult(success=False, error="python-pptx is not installed")
    except Exception as e:
        return ToolResult(success=False, error=str(e))


# ── Registration ──

def register_office_tools(tm, event_bus=None):
    import asyncio
    from aios.core.tool_manager import ToolContract
    from aios.core.permission_manager import PermissionLevel

    pdf_tools = [
        ToolContract(
            id="office.read_pdf", name="Read PDF",
            description="Extract text and metadata from a PDF file",
            parameters={
                "path": {"type": "string", "description": "Path to PDF file"},
            },
            returns={"path": {"type": "string"}, "pages": {"type": "array"}, "page_count": {"type": "integer"}},
            permission_level=PermissionLevel.READ, requires_confirmation=False, category="office",
            capabilities=["office.read_pdf"], tags=["office", "pdf", "read"],
        ),
        ToolContract(
            id="office.search_pdf", name="Search PDF",
            description="Search for text within a PDF file",
            parameters={
                "path": {"type": "string", "description": "Path to PDF file"},
                "query": {"type": "string", "description": "Text to search for"},
                "case_sensitive": {"type": "boolean", "description": "Case sensitive", "default": False},
            },
            returns={"path": {"type": "string"}, "query": {"type": "string"}, "results": {"type": "array"}, "matches": {"type": "integer"}},
            permission_level=PermissionLevel.READ, requires_confirmation=False, category="office",
            capabilities=["office.search_pdf"], tags=["office", "pdf", "search"],
        ),
        ToolContract(
            id="office.extract_pdf_metadata", name="Extract PDF Metadata",
            description="Extract metadata from a PDF file",
            parameters={
                "path": {"type": "string", "description": "Path to PDF file"},
            },
            returns={"path": {"type": "string"}, "page_count": {"type": "integer"}, "metadata": {"type": "object"}},
            permission_level=PermissionLevel.READ, requires_confirmation=False, category="office",
            capabilities=["office.extract_pdf_metadata"], tags=["office", "pdf", "metadata"],
        ),
        ToolContract(
            id="office.split_pdf", name="Split PDF",
            description="Split a PDF into separate files by page (requires confirmation)",
            parameters={
                "path": {"type": "string", "description": "Path to PDF file"},
                "output_dir": {"type": "string", "description": "Output directory", "required": False},
                "pages": {"type": "array", "description": "Specific page numbers to extract (1-based)", "default": [], "required": False},
            },
            returns={"source": {"type": "string"}, "files": {"type": "array"}, "count": {"type": "integer"}},
            permission_level=PermissionLevel.WORKSPACE, category="office",
            requires_confirmation=True,
            capabilities=["office.split_pdf"], tags=["office", "pdf", "split"],
        ),
        ToolContract(
            id="office.merge_pdf", name="Merge PDF",
            description="Merge multiple PDF files into one (requires confirmation)",
            parameters={
                "sources": {"type": "array", "description": "List of PDF file paths to merge"},
                "output": {"type": "string", "description": "Output PDF path"},
            },
            returns={"output": {"type": "string"}, "files": {"type": "array"}, "total_files": {"type": "integer"}},
            permission_level=PermissionLevel.WORKSPACE, category="office",
            requires_confirmation=True,
            capabilities=["office.merge_pdf"], tags=["office", "pdf", "merge"],
        ),
    ]

    docx_tools = [
        ToolContract(
            id="office.read_docx", name="Read DOCX",
            description="Read text content from a DOCX file with paragraph styles",
            parameters={
                "path": {"type": "string", "description": "Path to DOCX file"},
            },
            returns={"path": {"type": "string"}, "paragraphs": {"type": "array"}, "paragraph_count": {"type": "integer"}},
            permission_level=PermissionLevel.READ, requires_confirmation=False, category="office",
            capabilities=["office.read_docx"], tags=["office", "docx", "read"],
        ),
        ToolContract(
            id="office.write_docx", name="Write DOCX",
            description="Create a DOCX file with styled content (requires confirmation)",
            parameters={
                "path": {"type": "string", "description": "Path to output DOCX file"},
                "content": {"type": "array", "description": "Array of {text, style} objects"},
            },
            returns={"path": {"type": "string"}, "items_written": {"type": "integer"}},
            permission_level=PermissionLevel.WORKSPACE, category="office",
            requires_confirmation=True,
            capabilities=["office.write_docx"], tags=["office", "docx", "write"],
        ),
        ToolContract(
            id="office.extract_headings", name="Extract Headings",
            description="Extract heading structure from a DOCX file",
            parameters={
                "path": {"type": "string", "description": "Path to DOCX file"},
            },
            returns={"path": {"type": "string"}, "headings": {"type": "array"}, "count": {"type": "integer"}},
            permission_level=PermissionLevel.READ, requires_confirmation=False, category="office",
            capabilities=["office.extract_headings"], tags=["office", "docx", "headings"],
        ),
    ]

    xlsx_tools = [
        ToolContract(
            id="office.list_sheets", name="List Sheets",
            description="List all sheet names in an XLSX file",
            parameters={
                "path": {"type": "string", "description": "Path to XLSX file"},
            },
            returns={"path": {"type": "string"}, "sheets": {"type": "array"}, "count": {"type": "integer"}},
            permission_level=PermissionLevel.READ, requires_confirmation=False, category="office",
            capabilities=["office.list_sheets"], tags=["office", "xlsx", "sheets"],
        ),
        ToolContract(
            id="office.read_sheet", name="Read Sheet",
            description="Read data from a specific sheet in an XLSX file",
            parameters={
                "path": {"type": "string", "description": "Path to XLSX file"},
                "sheet": {"type": "string", "description": "Sheet name (omit for active sheet)", "required": False},
                "max_rows": {"type": "integer", "description": "Max rows to read", "default": 10000},
                "has_header": {"type": "boolean", "description": "First row is header", "default": True},
            },
            returns={"path": {"type": "string"}, "sheet": {"type": "string"}, "rows": {"type": "array"}, "row_count": {"type": "integer"}},
            permission_level=PermissionLevel.READ, requires_confirmation=False, category="office",
            capabilities=["office.read_sheet"], tags=["office", "xlsx", "read"],
        ),
        ToolContract(
            id="office.write_sheet", name="Write Sheet",
            description="Write data to an XLSX sheet (requires confirmation)",
            parameters={
                "path": {"type": "string", "description": "Path to XLSX file"},
                "sheet": {"type": "string", "description": "Sheet name", "default": "Sheet1"},
                "data": {"type": "array", "description": "Array of row objects or arrays"},
                "headers": {"type": "array", "description": "Column headers", "required": False},
            },
            returns={"path": {"type": "string"}, "sheet": {"type": "string"}, "rows_written": {"type": "integer"}},
            permission_level=PermissionLevel.WORKSPACE, category="office",
            requires_confirmation=True,
            capabilities=["office.write_sheet"], tags=["office", "xlsx", "write"],
        ),
    ]

    pptx_tools = [
        ToolContract(
            id="office.read_presentation", name="Read Presentation",
            description="Read slides and shapes from a PPTX file",
            parameters={
                "path": {"type": "string", "description": "Path to PPTX file"},
            },
            returns={"path": {"type": "string"}, "slides": {"type": "array"}, "slide_count": {"type": "integer"}},
            permission_level=PermissionLevel.READ, requires_confirmation=False, category="office",
            capabilities=["office.read_presentation"], tags=["office", "pptx", "read"],
        ),
        ToolContract(
            id="office.list_slides", name="List Slides",
            description="List all slides in a PPTX with titles",
            parameters={
                "path": {"type": "string", "description": "Path to PPTX file"},
            },
            returns={"path": {"type": "string"}, "slides": {"type": "array"}, "count": {"type": "integer"}},
            permission_level=PermissionLevel.READ, requires_confirmation=False, category="office",
            capabilities=["office.list_slides"], tags=["office", "pptx", "slides"],
        ),
        ToolContract(
            id="office.extract_notes", name="Extract Notes",
            description="Extract speaker notes from all slides in a PPTX",
            parameters={
                "path": {"type": "string", "description": "Path to PPTX file"},
            },
            returns={"path": {"type": "string"}, "notes": {"type": "array"}, "slides_with_notes": {"type": "integer"}},
            permission_level=PermissionLevel.READ, requires_confirmation=False, category="office",
            capabilities=["office.extract_notes"], tags=["office", "pptx", "notes"],
        ),
    ]

    all_tools = pdf_tools + docx_tools + xlsx_tools + pptx_tools

    pdf_handlers = [
        _read_pdf, _search_pdf, _extract_pdf_metadata,
        lambda p, eb=event_bus: _split_pdf(p, eb),
        lambda p, eb=event_bus: _merge_pdf(p, eb),
    ]
    docx_handlers = [
        _read_docx,
        lambda p, eb=event_bus: _write_docx(p, eb),
        _extract_headings,
    ]
    xlsx_handlers = [
        _list_sheets, _read_sheet,
        lambda p, eb=event_bus: _write_sheet(p, eb),
    ]
    pptx_handlers = [
        _read_presentation, _list_slides, _extract_notes,
    ]

    all_handlers = pdf_handlers + docx_handlers + xlsx_handlers + pptx_handlers

    for contract, handler in zip(all_tools, all_handlers):
        asyncio.create_task(tm.register_tool(contract, handler))
