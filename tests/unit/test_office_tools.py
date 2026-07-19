"""Tests for Office Document Toolkit (PDF, DOCX, XLSX, PPTX)."""

import asyncio
import json
from pathlib import Path

import pytest

from aios.core.tool_manager import ToolManager, ToolResult
from aios.core.permission_manager import PermissionManager
from aios.core.event_bus import EventBus
from aios.tools.office_tools import (
    register_office_tools,
    _read_pdf,
    _search_pdf,
    _extract_pdf_metadata,
    _split_pdf,
    _merge_pdf,
    _read_docx,
    _write_docx,
    _extract_headings,
    _list_sheets,
    _read_sheet,
    _write_sheet,
    _read_presentation,
    _list_slides,
    _extract_notes,
)


@pytest.fixture
def pm():
    return PermissionManager()


@pytest.fixture
def tm(pm):
    return ToolManager(pm)


@pytest.fixture
async def eb():
    bus = EventBus(max_retries=1, retry_delay=0.01)
    await bus.start()
    yield bus
    await bus.stop()


# ── Fixtures: Create test documents ──


@pytest.fixture
def sample_pdf(tmp_path):
    import fitz
    f = tmp_path / "test.pdf"
    doc = fitz.open()
    doc.new_page()
    doc[0].insert_text((50, 50), "Hello PDF World")
    page2 = doc.new_page()
    page2.insert_text((50, 50), "Second page content")
    doc.save(str(f))
    doc.close()
    return f


@pytest.fixture
def sample_docx(tmp_path):
    from docx import Document
    f = tmp_path / "test.docx"
    doc = Document()
    doc.add_heading("Title Here", 0)
    doc.add_paragraph("First paragraph.")
    doc.add_heading("Section 1", 1)
    doc.add_paragraph("Section body.")
    doc.add_heading("Section 2", 2)
    doc.add_paragraph("More text.")
    doc.save(str(f))
    return f


@pytest.fixture
def sample_xlsx(tmp_path):
    from openpyxl import Workbook
    f = tmp_path / "test.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Age", "City"])
    ws.append(["Alice", 30, "NYC"])
    ws.append(["Bob", 25, "LA"])
    ws.append(["Carol", 35, "Chicago"])
    ws2 = wb.create_sheet("DataSheet")
    ws2.append(["X", "Y"])
    ws2.append([1, 10])
    ws2.append([2, 20])
    wb.save(str(f))
    wb.close()
    return f


@pytest.fixture
def sample_pptx(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    f = tmp_path / "test.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Slide One"
    slide1.notes_slide.notes_text_frame.text = "Note for slide 1"
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Slide Two"
    prs.save(str(f))
    return f


# ═══════════════════════════════════════════════════════════════════
# PDF Tests
# ═══════════════════════════════════════════════════════════════════


@pytest.mark.asyncio
async def test_read_pdf(sample_pdf):
    result = await _read_pdf({"path": str(sample_pdf)})
    assert result.success
    assert result.data["page_count"] == 2
    assert "Hello PDF World" in result.data["pages"][0]["text"]
    assert "Second page" in result.data["pages"][1]["text"]


@pytest.mark.asyncio
async def test_read_pdf_not_found():
    result = await _read_pdf({"path": "/nonexistent.pdf"})
    assert not result.success


@pytest.mark.asyncio
async def test_search_pdf(sample_pdf):
    result = await _search_pdf({"path": str(sample_pdf), "query": "Hello"})
    assert result.success
    assert result.data["matches"] >= 1


@pytest.mark.asyncio
async def test_search_pdf_no_match(sample_pdf):
    result = await _search_pdf({"path": str(sample_pdf), "query": "ZZZZ"})
    assert result.success
    assert result.data["matches"] == 0


@pytest.mark.asyncio
async def test_search_pdf_no_query(sample_pdf):
    result = await _search_pdf({"path": str(sample_pdf), "query": ""})
    assert not result.success


@pytest.mark.asyncio
async def test_extract_pdf_metadata(sample_pdf):
    result = await _extract_pdf_metadata({"path": str(sample_pdf)})
    assert result.success
    assert result.data["page_count"] == 2
    assert "metadata" in result.data


@pytest.mark.asyncio
async def test_split_pdf(sample_pdf, eb):
    result = await _split_pdf({"path": str(sample_pdf)}, eb)
    assert result.success
    assert result.data["count"] == 2


@pytest.mark.asyncio
async def test_split_pdf_specific_pages(sample_pdf, eb):
    result = await _split_pdf({"path": str(sample_pdf), "pages": [1]}, eb)
    assert result.success
    assert result.data["count"] == 1


@pytest.mark.asyncio
async def test_merge_pdf(sample_pdf, eb):
    out = sample_pdf.parent / "merged.pdf"
    result = await _merge_pdf({"sources": [str(sample_pdf), str(sample_pdf)], "output": str(out)}, eb)
    assert result.success
    assert result.data["total_files"] == 2
    assert out.exists()


@pytest.mark.asyncio
async def test_merge_pdf_no_sources(eb):
    result = await _merge_pdf({"sources": [], "output": "/dev/null/out.pdf"}, eb)
    assert not result.success


# ═══════════════════════════════════════════════════════════════════
# DOCX Tests
# ═══════════════════════════════════════════════════════════════════


@pytest.mark.asyncio
async def test_read_docx(sample_docx):
    result = await _read_docx({"path": str(sample_docx)})
    assert result.success
    assert result.data["paragraph_count"] >= 5


@pytest.mark.asyncio
async def test_read_docx_content(sample_docx):
    result = await _read_docx({"path": str(sample_docx)})
    texts = [p["text"] for p in result.data["paragraphs"]]
    assert any("Title Here" in t for t in texts)
    assert any("First paragraph" in t for t in texts)
    assert any("Section 1" in t for t in texts)


@pytest.mark.asyncio
async def test_read_docx_headings(sample_docx):
    result = await _read_docx({"path": str(sample_docx)})
    headings = result.data["headings"]
    assert len(headings) >= 2


@pytest.mark.asyncio
async def test_read_docx_not_found():
    result = await _read_docx({"path": "/nonexistent.docx"})
    assert not result.success


@pytest.mark.asyncio
async def test_write_docx(tmp_path, eb):
    out = tmp_path / "created.docx"
    content = [
        {"text": "My Document", "style": "title"},
        {"text": "Hello world", "style": "paragraph"},
        {"text": "Chapter 1", "style": "heading1"},
    ]
    result = await _write_docx({"path": str(out), "content": content}, eb)
    assert result.success
    assert result.data["items_written"] == 3
    assert out.exists()


@pytest.mark.asyncio
async def test_extract_headings(sample_docx):
    result = await _extract_headings({"path": str(sample_docx)})
    assert result.success
    assert result.data["count"] >= 2
    texts = [h["text"] for h in result.data["headings"]]
    assert "Section 1" in texts


@pytest.mark.asyncio
async def test_extract_headings_not_found():
    result = await _extract_headings({"path": "/nope.docx"})
    assert not result.success


# ═══════════════════════════════════════════════════════════════════
# XLSX Tests
# ═══════════════════════════════════════════════════════════════════


@pytest.mark.asyncio
async def test_list_sheets(sample_xlsx):
    result = await _list_sheets({"path": str(sample_xlsx)})
    assert result.success
    assert result.data["count"] == 2
    names = [s["name"] for s in result.data["sheets"]]
    assert "Sheet1" in names
    assert "DataSheet" in names


@pytest.mark.asyncio
async def test_read_sheet(sample_xlsx):
    result = await _read_sheet({"path": str(sample_xlsx)})
    assert result.success
    assert result.data["row_count"] == 3


@pytest.mark.asyncio
async def test_read_sheet_by_name(sample_xlsx):
    result = await _read_sheet({"path": str(sample_xlsx), "sheet": "DataSheet"})
    assert result.success
    assert result.data["row_count"] == 2


@pytest.mark.asyncio
async def test_read_sheet_not_found(sample_xlsx):
    result = await _read_sheet({"path": str(sample_xlsx), "sheet": "NoSuchSheet"})
    assert not result.success


@pytest.mark.asyncio
async def test_read_sheet_no_file():
    result = await _read_sheet({"path": "/nonexistent.xlsx"})
    assert not result.success


@pytest.mark.asyncio
async def test_write_sheet(tmp_path, eb):
    out = tmp_path / "new.xlsx"
    data = [{"name": "Alice", "age": 30}, {"name": "Bob", "age": 25}]
    result = await _write_sheet({"path": str(out), "data": data, "headers": ["name", "age"]}, eb)
    assert result.success
    assert result.data["rows_written"] == 2
    assert out.exists()


@pytest.mark.asyncio
async def test_write_sheet_no_data(tmp_path, eb):
    result = await _write_sheet({"path": str(tmp_path / "empty.xlsx"), "data": []}, eb)
    assert not result.success


# ═══════════════════════════════════════════════════════════════════
# PPTX Tests
# ═══════════════════════════════════════════════════════════════════


@pytest.mark.asyncio
async def test_read_presentation(sample_pptx):
    result = await _read_presentation({"path": str(sample_pptx)})
    assert result.success
    assert result.data["slide_count"] == 2
    assert result.data["slides"][0]["title"] == "Slide One"
    assert result.data["slides"][1]["title"] == "Slide Two"


@pytest.mark.asyncio
async def test_read_presentation_notes(sample_pptx):
    result = await _read_presentation({"path": str(sample_pptx)})
    assert result.data["slides"][0].get("notes") == "Note for slide 1"


@pytest.mark.asyncio
async def test_list_slides(sample_pptx):
    result = await _list_slides({"path": str(sample_pptx)})
    assert result.success
    assert result.data["count"] == 2
    assert result.data["slides"][0]["title"] == "Slide One"


@pytest.mark.asyncio
async def test_extract_notes(sample_pptx):
    result = await _extract_notes({"path": str(sample_pptx)})
    assert result.success
    assert result.data["slides_with_notes"] == 1
    assert result.data["notes"][0]["notes"] == "Note for slide 1"
    assert not result.data["notes"][1]["has_notes"]


@pytest.mark.asyncio
async def test_read_presentation_not_found():
    result = await _read_presentation({"path": "/nope.pptx"})
    assert not result.success


# ═══════════════════════════════════════════════════════════════════
# Corrupted Document Handling
# ═══════════════════════════════════════════════════════════════════


@pytest.mark.asyncio
async def test_read_corrupted_pdf(tmp_path):
    f = tmp_path / "bad.pdf"
    f.write_text("this is not a pdf")
    result = await _read_pdf({"path": str(f)})
    assert not result.success


@pytest.mark.asyncio
async def test_read_corrupted_docx(tmp_path):
    f = tmp_path / "bad.docx"
    f.write_text("not a docx")
    result = await _read_docx({"path": str(f)})
    assert not result.success


@pytest.mark.asyncio
async def test_read_corrupted_xlsx(tmp_path):
    f = tmp_path / "bad.xlsx"
    f.write_text("not an xlsx")
    result = await _read_sheet({"path": str(f)})
    assert not result.success


@pytest.mark.asyncio
async def test_read_corrupted_pptx(tmp_path):
    f = tmp_path / "bad.pptx"
    f.write_text("not a pptx")
    result = await _read_presentation({"path": str(f)})
    assert not result.success


# ═══════════════════════════════════════════════════════════════════
# Registration Tests
# ═══════════════════════════════════════════════════════════════════


@pytest.mark.asyncio
async def test_register_office_tools(tm, eb):
    register_office_tools(tm, eb)
    await asyncio.sleep(0.05)
    all_tools = await tm.list_tools()
    office_tools = [t for t in all_tools if t.category == "office"]
    ids = {t.id for t in office_tools}

    expected = {
        "office.read_pdf", "office.search_pdf", "office.extract_pdf_metadata",
        "office.split_pdf", "office.merge_pdf",
        "office.read_docx", "office.write_docx", "office.extract_headings",
        "office.list_sheets", "office.read_sheet", "office.write_sheet",
        "office.read_presentation", "office.list_slides", "office.extract_notes",
    }
    assert ids == expected, f"Missing: {expected - ids}, Extra: {ids - expected}"
    assert len(office_tools) == 14


@pytest.mark.asyncio
async def test_office_tool_ids_unique(tm, eb):
    register_office_tools(tm, eb)
    await asyncio.sleep(0.05)
    all_tools = await tm.list_tools()
    office_tools = [t for t in all_tools if t.category == "office"]
    ids = [t.id for t in office_tools]
    assert len(ids) == len(set(ids)), "Duplicate tool IDs found"


@pytest.mark.asyncio
async def test_office_tools_have_permission_levels(tm, eb):
    register_office_tools(tm, eb)
    await asyncio.sleep(0.05)
    all_tools = await tm.list_tools()
    office_tools = [t for t in all_tools if t.category == "office"]
    for t in office_tools:
        assert t.permission_level is not None
        assert 0 <= int(t.permission_level) <= 3


@pytest.mark.asyncio
async def test_write_office_tools_require_confirmation(tm, eb):
    register_office_tools(tm, eb)
    await asyncio.sleep(0.05)
    write_ids = {
        "office.split_pdf", "office.merge_pdf",
        "office.write_docx", "office.write_sheet",
    }
    for tid in write_ids:
        contract = await tm.get_tool(tid)
        assert contract is not None, f"{tid} not registered"
        assert contract.requires_confirmation, f"{tid} should require confirmation"


@pytest.mark.asyncio
async def test_read_office_tools_no_confirmation(tm, eb):
    register_office_tools(tm, eb)
    await asyncio.sleep(0.05)
    read_ids = {
        "office.read_pdf", "office.search_pdf", "office.extract_pdf_metadata",
        "office.read_docx", "office.extract_headings",
        "office.list_sheets", "office.read_sheet",
        "office.read_presentation", "office.list_slides", "office.extract_notes",
    }
    for tid in read_ids:
        contract = await tm.get_tool(tid)
        assert contract is not None, f"{tid} not registered"
        assert not contract.requires_confirmation, f"{tid} should not require confirmation"


# ═══════════════════════════════════════════════════════════════════
# Large Document Tests
# ═══════════════════════════════════════════════════════════════════


@pytest.mark.asyncio
async def test_large_pdf(tmp_path):
    import fitz
    f = tmp_path / "large.pdf"
    doc = fitz.open()
    for _ in range(50):
        p = doc.new_page()
        p.insert_text((50, 50), "Page content here. " * 20)
    doc.save(str(f))
    doc.close()

    result = await _read_pdf({"path": str(f)})
    assert result.success
    assert result.data["page_count"] == 50


@pytest.mark.asyncio
async def test_large_xlsx(tmp_path, eb):
    from openpyxl import Workbook
    f = tmp_path / "large.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.append(["col"] * 20)
    for i in range(500):
        ws.append([f"val_{i}_{j}" for j in range(20)])
    wb.save(str(f))
    wb.close()

    result = await _read_sheet({"path": str(f), "max_rows": 100, "has_header": False})
    assert result.success
    assert result.data["row_count"] == 100


# ═══════════════════════════════════════════════════════════════════
# XLSX List Sheets with no file
# ═══════════════════════════════════════════════════════════════════


@pytest.mark.asyncio
async def test_list_sheets_not_found():
    result = await _list_sheets({"path": "/nope.xlsx"})
    assert not result.success


@pytest.mark.asyncio
async def test_extract_notes_not_found():
    result = await _extract_notes({"path": "/nope.pptx"})
    assert not result.success


@pytest.mark.asyncio
async def test_list_slides_not_found():
    result = await _list_slides({"path": "/nope.pptx"})
    assert not result.success
